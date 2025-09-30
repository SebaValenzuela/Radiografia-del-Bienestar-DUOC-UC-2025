import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Pt
import os
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import numpy as np
import matplotlib.ticker as mticker
import pandas as pd

# Estilo global para los gráficos
plt.style.use('ggplot')
plt.rcParams.update({
    'font.size': 12,
    'figure.dpi': 150
})

def crear_grafico_pie(df, col_respuestas, col_total, output_path):
    responded = df[col_respuestas].sum()
    not_responded = max(df[col_total].sum() - responded, 0)
    sizes = [responded, not_responded]
    labels_base = ['Encuestas respondidas', 'Encuestas no respondidas']
    
    # --- Crear labels con porcentaje y valor ---
    total = sum(sizes)
    if total == 0:
        labels = [f"{label}\n0 (0%)" for label in labels_base]
    else:
        labels = [f"{label}\n{size} ({size/total*100:.2f}%)"
                for label, size in zip(labels_base, sizes)]
    
    colors = ['#000000', '#F8B416']  # negro y amarillo
    explode = (0.05, 0.05)

    fig, ax = plt.subplots(figsize=(5,5))

    wedges, texts = ax.pie(
        sizes,
        labels=labels,
        startangle=90,
        colors=colors,
        explode=explode,
        wedgeprops={'edgecolor':'white', 'linewidth':1.5},
        textprops={'fontsize':12},
    )

    centre_circle = plt.Circle((0,-0.05), 0.60, fc='white')
    ax.add_artist(centre_circle)

    ax.axis('equal')
    plt.tight_layout()
    plt.savefig(output_path, transparent=True, bbox_inches='tight')
    print("Gráfico de pastel guardado en:", output_path)
    plt.close()

def crear_grafico_barras(df, col_categoria, col_valor, output_path, fontsize=12):
    fig, ax = plt.subplots(figsize=(6,4))

    df_plot = df.copy()

    if col_valor not in df_plot.columns:
        raise ValueError(f"Columna '{col_valor}' no encontrada en el DataFrame")

    s = df_plot[col_valor]

    s_str = s.astype(str).str.replace(r'\s+', '', regex=True)
    s_clean = s_str.str.rstrip('%').str.replace(',', '.', regex=False)

    s_num = pd.to_numeric(s_clean, errors='coerce')

    if pd.notna(s_num.max()) and s_num.max() <= 1.1:
        s_num = s_num * 100

    s_num = s_num.fillna(0)

    df_plot['_plot_val'] = s_num

    x = np.arange(len(df_plot))
    bars = ax.bar(
        x,
        df_plot['_plot_val'],
        color='#1f77b4',      # azul por defecto (puedes poner otro código HEX)
        edgecolor='black',
        width=0.4
    )

    ymax = max(100, df_plot['_plot_val'].max() * 1.1)
    ax.set_ylim(0, 20)

    # mostrar eje como porcentaje (0..100 -> "0%..100%")
    ax.yaxis.set_major_formatter(mticker.PercentFormatter(xmax=100, decimals=0))

    # configuración visual similar a la original
    ax.tick_params(left=False, bottom=False)
    ax.set_xticks(x)
    ax.set_xticklabels(df_plot[col_categoria], rotation=45, ha='right', fontsize=fontsize)

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.spines['bottom'].set_visible(False)

    grid_color = '#CCCCCC'
    grid_linestyle = '-'
    ax.yaxis.grid(True, color=grid_color, linestyle=grid_linestyle)
    ax.xaxis.grid(False)

    plt.tight_layout()
    plt.savefig(output_path, transparent=True, bbox_inches='tight')
    plt.close()

def rellenar_tabla(slide, placeholder_name, df):
    for shape in slide.shapes:
        if shape.has_text_frame and placeholder_name in shape.text:
            rows, cols = df.shape[0] + 1, df.shape[1]
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            table = slide.shapes.add_table(
                rows, cols, left, top, int(width), int(height)
            ).table

            # ajustar anchos y altos
            for col in table.columns:
                col.width = int((width*0.8) // cols)
            for row in table.rows:
                row.height = int((height*0.6) // rows)

            # encabezado
            for j, col_name in enumerate(df.columns):
                cell = table.cell(0, j)
                cell.text = str(col_name)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(142, 209, 252)
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0]
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 0, 0)

            # contenido
            for i in range(df.shape[0]):
                for j in range(df.shape[1]):
                    cell = table.cell(i+1, j)
                    cell.text = str(df.iloc[i, j])
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0]
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor(50, 50, 50)

                    # Reglas de colores SOLO para tabla_avance_sedes
                    if placeholder_name == "tabla_avance_sedes" and df.columns[j] == "Cantidad de respuestas":
                        if "N referencial" in df.columns:
                            respuestas = df.iloc[i][j]
                            n_ref = df.iloc[i]["N referencial"]
                            if n_ref > 0:
                                ratio = respuestas / n_ref
                                if ratio < 0.75:
                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = RGBColor(255, 165, 0)  # naranja
                                else:
                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = RGBColor(0, 176, 80)  # verde

                    # Otras tablas mantienen el diseño original
                    elif df.columns[j] == "% de avance respecto a total":
                        # lógica top3/bottom3 si aplica
                        pass

            shape.text = ""  # limpiar placeholder
            break


def rellenar_grafico(slide, placeholder_name, image_path):
    for shape in slide.shapes:
        if shape.has_text_frame and placeholder_name in shape.text:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            slide.shapes.add_picture(image_path, left, top, width, height)
            shape.text = ""  # eliminar placeholder
            break

def rellenar_tabla_parcial(prs, df, base_placeholder="tabla_resumen_escuela_en_cada_sede", max_rows_por_slide=10):
    total_rows = df.shape[0]
    num_splits = (total_rows + max_rows_por_slide - 1) // max_rows_por_slide  # cuántas slides necesitas

    for i in range(num_splits):
        start_row = i * max_rows_por_slide
        end_row = min(start_row + max_rows_por_slide, total_rows)
        df_parcial = df.iloc[start_row:end_row]

        placeholder_name = f"{base_placeholder}_{i+1}"  # ejemplo: tabla_resumen_escuela_en_cada_sede_1
        slide = None

        # buscar slide que tenga ese placeholder
        for s in prs.slides:
            for shape in s.shapes:
                if shape.has_text_frame and placeholder_name in shape.text:
                    slide = s
                    break
            if slide:
                break

        if not slide:
            continue

        # rellenar tabla igual que antes
        for shape in slide.shapes:
            if shape.has_text_frame and placeholder_name in shape.text:
                rows, cols = df_parcial.shape[0] + 1, df_parcial.shape[1]
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table

                # ajustar ancho/alto
                for col in table.columns:
                    col.width = int((width*0.8) // cols)
                for row in table.rows:
                    row.height = int((height*0.6) // rows)

                # encabezado
                for j, col_name in enumerate(df_parcial.columns):
                    cell = table.cell(0, j)
                    cell.text = str(col_name)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(142, 209, 252)
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0]
                    run.font.size = Pt(9)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 0, 0)

                # top3/bottom3
                if "%" in df_parcial.columns:
                    valores = df_parcial["%"]
                    top3_idx = valores.nlargest(3).index
                    bottom3_idx = valores.nsmallest(3).index
                else:
                    top3_idx, bottom3_idx = [], []

                # contenido
                for r in range(df_parcial.shape[0]):
                    for c in range(df_parcial.shape[1]):
                        cell = table.cell(r+1, c)
                        cell.text = str(df_parcial.iloc[r, c])
                        p = cell.text_frame.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        run = p.runs[0]
                        run.font.size = Pt(8)
                        run.font.color.rgb = RGBColor(50, 50, 50)

                        if df_parcial.columns[c] == "% de avance respecto a total":
                            if df_parcial.index[r] in top3_idx:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(11, 241, 30)  # verde
                            elif df_parcial.index[r] in bottom3_idx:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(229, 239, 29)  # amarillo

                shape.text = ""  # limpiar placeholder
                break


def generar_presentacion(template_path, output_path, resumen_sede, resumen_escuela, resumen_escuela_y_sede):
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs = Presentation(template_path)

    # --- Crear gráficos ---
    crear_grafico_pie(resumen_sede, 'Cantidad de respuestas', 'Cantidad de estudiantes', "grafico_avance_global.png")
    crear_grafico_barras(resumen_sede, 'SEDE', '% de avance respecto a total', "grafico_avance_sedes.png", fontsize=10)
    crear_grafico_barras(resumen_escuela, 'ESCUELA', '% de avance respecto a total', "grafico_avance_escuelas.png", fontsize=8)

    # --- Buscar slide template para la tabla larga ---
    template_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and "tabla_resumen_escuela_en_cada_sede" in shape.text:
                template_slide = slide
                break
        if template_slide:
            break

    # --- Reemplazar placeholders de gráficos y tablas ---
    for slide in prs.slides:
        rellenar_grafico(slide, "grafico_avance_global", "grafico_avance_global.png")
        rellenar_tabla(slide, "tabla_avance_sedes", resumen_sede)
        rellenar_grafico(slide, "grafico_avance_sedes", "grafico_avance_sedes.png")
        rellenar_tabla(slide, "tabla_avance_escuelas", resumen_escuela)
        rellenar_grafico(slide, "grafico_avance_escuelas", "grafico_avance_escuelas.png")
        rellenar_tabla_parcial(prs, resumen_escuela_y_sede, max_rows_por_slide=9)

    prs.save(output_path)
    print(f"Presentación generada en '{output_path}'")

